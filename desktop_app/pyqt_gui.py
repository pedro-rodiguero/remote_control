from PyQt5.QtWidgets import QApplication, QMainWindow, QLabel, QVBoxLayout, QWidget, QPushButton, QFileDialog, QMessageBox, QComboBox, QDialog, QFormLayout, QHBoxLayout, QGridLayout
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QPixmap
import qrcode
import socket
import sys
from pptx import Presentation
from pptx.util import Inches
from PIL import Image

class DragDropLabel(QLabel):
    def __init__(self, parent=None, app_instance=None):
        super().__init__(parent)
        self.app_instance = app_instance
        self.setAcceptDrops(True)
        self.setText("Drag and drop files here")
        self.setAlignment(Qt.AlignCenter)
        self.setStyleSheet("border: 2px dashed #aaa;")

    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dropEvent(self, event):
        for url in event.mimeData().urls():
            file_path = url.toLocalFile()
            self.app_instance.upload_file(file_path)

def get_local_ip():
    s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
    try:
        # doesn't even have to be reachable
        s.connect(('10.254.254.254', 1))
        ip = s.getsockname()[0]
    except Exception:
        ip = '127.0.0.1'
    finally:
        s.close()
    return ip

def generate_qr_code(url):
    qr = qrcode.QRCode(
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L,
        box_size=10,
        border=4,
    )
    qr.add_data(url)
    qr.make(fit=True)
    img = qr.make_image(fill='black', back_color='white')
    qr_path = "qr_code.png"
    img.save(qr_path)
    return qr_path

class App(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Presentation Software")
        self.setGeometry(100, 100, 600, 400)
        self.presentation_uploaded = False  # Flag to track if a presentation is uploaded
        self.current_slide_index = 0
        self.slides = []
        self.initUI()

    def initUI(self):
        central_widget = QWidget()
        self.setCentralWidget(central_widget)

        self.layout = QVBoxLayout()

        self.home_button = QPushButton("Home", self)
        self.home_button.clicked.connect(self.show_home)
        self.layout.addWidget(self.home_button)

        self.start_button = QPushButton("Start Casting", self)
        self.start_button.clicked.connect(self.start_casting)
        self.layout.addWidget(self.start_button)

        self.stop_button = QPushButton("Stop Casting", self)
        self.stop_button.clicked.connect(self.stop_casting)
        self.layout.addWidget(self.stop_button)

        self.select_screen_button = QPushButton("Select Output Screen", self)
        self.select_screen_button.clicked.connect(self.open_screen_selection)
        self.layout.addWidget(self.select_screen_button)

        self.qr_button = QPushButton("Generate QR Code", self)
        self.qr_button.clicked.connect(self.generate_qr_code)
        self.layout.addWidget(self.qr_button)

        self.upload_button = QPushButton("Upload Presentation", self)
        self.upload_button.clicked.connect(self.upload_presentation)
        self.layout.addWidget(self.upload_button)

        self.drag_drop_label = DragDropLabel(self, app_instance=self)
        self.layout.addWidget(self.drag_drop_label)

        self.grid_layout = QGridLayout()
        self.layout.addLayout(self.grid_layout)

        if self.presentation_uploaded:
            self.init_grid_layout()
        else:
            self.drag_drop_label.show()

        central_widget.setLayout(self.layout)

    def init_grid_layout(self):
        self.clear_grid_layout()

        # Add labels for the grid layout
        self.current_slide_label = QLabel("Current Slide")
        self.current_slide_label.setStyleSheet("font-weight: bold;")
        self.grid_layout.addWidget(self.current_slide_label, 0, 0, 1, 2)

        self.current_slide_image = QLabel()
        self.current_slide_image.setStyleSheet("border: 1px solid black;")
        self.grid_layout.addWidget(self.current_slide_image, 1, 0, 1, 2)

        self.prev_slide_button = QPushButton("Previous Slide")
        self.prev_slide_button.clicked.connect(self.prev_slide)
        self.grid_layout.addWidget(self.prev_slide_button, 2, 0)

        self.next_slide_button = QPushButton("Next Slide")
        self.next_slide_button.clicked.connect(self.next_slide)
        self.grid_layout.addWidget(self.next_slide_button, 2, 1)

    def start_casting(self):
        # Implement start casting functionality
        pass

    def stop_casting(self):
        # Implement stop casting functionality
        pass

    def open_screen_selection(self):
        dialog = QDialog(self)
        dialog.setWindowTitle("Select Output Screen")
        layout = QFormLayout(dialog)

        self.screen_type = QComboBox(dialog)
        self.screen_type.addItems(["Physical"])
        layout.addRow("Select Screen Type:", self.screen_type)

        self.output_screen = QComboBox(dialog)
        self.output_screen.addItems(["Screen 1", "Screen 2", "Screen 3"])
        layout.addRow("Select Output Screen:", self.output_screen)

        save_button = QPushButton("Save", dialog)
        save_button.clicked.connect(dialog.accept)
        layout.addWidget(save_button)

        dialog.exec_()

    def generate_qr_code(self):
        ip = get_local_ip()
        url = f"http://{ip}:5000"
        qr_path = generate_qr_code(url)
        self.show_qr_code(qr_path)

    def show_qr_code(self, qr_path):
        qr_window = QDialog(self)
        qr_window.setWindowTitle("QR Code")
        qr_window.setGeometry(100, 100, 300, 300)

        try:
            pixmap = QPixmap(qr_path)
            label = QLabel(qr_window)
            label.setPixmap(pixmap)
            label.setAlignment(Qt.AlignCenter)
            layout = QVBoxLayout()
            layout.addWidget(label)
            qr_window.setLayout(layout)
            qr_window.exec_()
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to load QR code image: {e}")

    def upload_presentation(self):
        file_path = QFileDialog.getOpenFileName(self, "Select Presentation", "", "PDF files (*.pdf);;PowerPoint files (*.pptx);;All files (*.*)")[0]
        if file_path:
            self.upload_file(file_path)
            self.presentation_uploaded = True
            self.update_ui_after_upload()

    def upload_file(self, file_path):
        try:
            if file_path.endswith('.pptx'):
                self.load_pptx(file_path)
            else:
                # Implement the logic to upload the file to the server
                pass
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to upload file: {e}")

    def load_pptx(self, file_path):
        presentation = Presentation(file_path)
        self.slides = []
        for slide in presentation.slides:
            image = slide.shapes.add_picture(file_path, Inches(0), Inches(0))
            self.slides.append(image)
        self.show_slide(0)

    def show_slide(self, index):
        if 0 <= index < len(self.slides):
            self.current_slide_index = index
            slide_image = self.slides[index]
            pixmap = QPixmap(slide_image)
            self.current_slide_image.setPixmap(pixmap)

    def update_ui_after_upload(self):
        self.select_screen_button.hide()
        self.upload_button.hide()
        self.drag_drop_label.hide()
        self.init_grid_layout()

    def show_home(self):
        self.select_screen_button.show()
        self.upload_button.show()
        self.drag_drop_label.show()
        self.presentation_uploaded = False
        self.clear_grid_layout()

    def clear_grid_layout(self):
        while self.grid_layout.count():
            item = self.grid_layout.takeAt(0)
            widget = item.widget()
            if widget is not None:
                widget.deleteLater()

    def prev_slide(self):
        self.show_slide(self.current_slide_index - 1)

    def next_slide(self):
        self.show_slide(self.current_slide_index + 1)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    main_window = App()
    main_window.show()
    sys.exit(app.exec_())