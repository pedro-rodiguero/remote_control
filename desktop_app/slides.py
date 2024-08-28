from pptx import Presentation
from pdf2image import convert_from_path
from PIL import Image

def load_pptx(file_path):
    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text)
        # Save slide as image
        slide_image_path = f"slide_{len(slides)}.png"
        slide.shapes._spTree.write(slide_image_path)
        slides.append(slide_image_path)
    return slides

def load_pdf(file_path):
    images = convert_from_path(file_path)
    slides = []
    for i, image in enumerate(images):
        slide_image_path = f"slide_{i}.png"
        image.save(slide_image_path, "PNG")
        slides.append(slide_image_path)
    return slides

def load_images(file_paths):
    return file_paths

def load_video(file_path):
    return [file_path]