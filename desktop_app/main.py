import tkinter as tk
from tkinter import messagebox, Toplevel, Label, filedialog
import pyautogui  # type: ignore
import qrcode  # type: ignore
import socket
import os
from PIL import Image, ImageTk
from flask import Flask, jsonify
from flask_cors import CORS
import threading
from pptx import Presentation  # type: ignore
from pdf2image import convert_from_path  # type: ignore
from ffpyplayer.player import MediaPlayer  # type: ignore
import pynetworktables as ndi  # Import the NDI library

app = Flask(__name__)
CORS(app)

class App:
    def __init__(self, root):
        self.root = root
        self.root.title("Remote Control Application")
        self.root.geometry("800x600")

        self.slides = []
        self.current_slide = 0
        self.player = None

        self.create_widgets()

    def create_widgets(self):
        self.load_button = tk.Button(
            self.root, text="Load Slides", command=self.load_slides
        )
        self.load_button.pack(pady=10)

        self.next_button = tk.Button(
            self.root, text="Next Slide", command=self.next_slide
        )
        self.next_button.pack(pady=10)

        self.prev_button = tk.Button(
            self.root, text="Previous Slide", command=self.prev_slide
        )
        self.prev_button.pack(pady=10)

        self.qr_button = tk.Button(
            self.root, text="Generate QR Code", command=self.generate_qr
        )
        self.qr_button.pack(pady=10)

        self.screen_button = tk.Button(
            self.root, text="Show Screen", command=self.show_screen
        )
        self.screen_button.pack(pady=10)

        self.slide_label = Label(self.root)
        self.slide_label.pack(pady=10)

    def load_slides(self):
        file_path = filedialog.askopenfilename(
            filetypes=[("PowerPoint files", "*.pptx"), ("PDF files", "*.pdf"), ("Image files", "*.png;*.jpg;*.jpeg"), ("Video files", "*.mp4;*.avi;*.mov")]
        )
        if file_path:
            if file_path.endswith(".pptx"):
                self.load_pptx(file_path)
            elif file_path.endswith(".pdf"):
                self.load_pdf(file_path)
            elif file_path.endswith((".mp4", ".avi", ".mov")):
                self.load_video(file_path)
            else:
                self.load_images([file_path])

    def load_pptx(self, file_path):
        prs = Presentation(file_path)
        self.slides = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        print(run.text)
            # Save slide as image
            slide_image_path = f"slide_{len(self.slides)}.png"
            slide.shapes._spTree.write(slide_image_path)
            self.slides.append(slide_image_path)
        self.show_slide(0)

    def load_pdf(self, file_path):
        images = convert_from_path(file_path)
        self.slides = []
        for i, image in enumerate(images):
            slide_image_path = f"slide_{i}.png"
            image.save(slide_image_path, "PNG")
            self.slides.append(slide_image_path)
        self.show_slide(0)

    def load_images(self, file_paths):
        self.slides = file_paths
        self.show_slide(0)

    def load_video(self, file_path):
        self.slides = [file_path]
        self.show_slide(0)

    def show_slide(self, index):
        if 0 <= index < len(self.slides):
            self.current_slide = index
            if self.player:
                self.player.close_player()
                self.player = None

            if self.slides[index].endswith((".mp4", ".avi", ".mov")):
                self.player = MediaPlayer(self.slides[index])
                self.play_video()
            else:
                img = Image.open(self.slides[index])
                img = img.resize((800, 600), Image.LANCZOS)
                img = ImageTk.PhotoImage(img)
                self.slide_label.config(image=img)
                self.slide_label.image = img

    def play_video(self):
        def update_frame():
            frame, val = self.player.get_frame()
            if val == 'eof':
                self.next_slide()
                return
            if frame is not None:
                img, t = frame
                img = Image.fromarray(img)
                img = img.resize((800, 600), Image.LANCZOS)
                img = ImageTk.PhotoImage(img)
                self.slide_label.config(image=img)
                self.slide_label.image = img
            self.root.after(10, update_frame)

        update_frame()

    def next_slide(self):
        if self.current_slide < len(self.slides) - 1:
            self.show_slide(self.current_slide + 1)

    def prev_slide(self):
        if self.current_slide > 0:
            self.show_slide(self.current_slide - 1)

    def generate_qr(self):
        ip = self.get_local_ip()
        url = f"http://{ip}:5000"
        print(f"Generated URL: {url}")  # Debug statement
        qr = qrcode.make(url)
        qr_path = os.path.join(os.getcwd(), "qrcode.png")
        qr.save(qr_path)
        self.show_qr_code(qr_path)

    def show_qr_code(self, qr_path):
        qr_window = Toplevel(self.root)
        qr_window.title("QR Code")
        qr_window.geometry("300x300")

        try:
            img = Image.open(qr_path)
            img = img.resize((250, 250), Image.LANCZOS)
            img = ImageTk.PhotoImage(img)

            label = Label(qr_window, image=img)
            label.image = img  # Keep a reference to avoid garbage collection
            label.pack(pady=10)
        except Exception as e:
            messagebox.showerror("Error", f"Failed to load QR code image: {e}")

    def show_screen(self):
        # Start NDI stream
        if not ndi.initialize():
            print("Cannot run NDI.")
            return

        # Create a sender
        sender = ndi.send_create()
        if not sender:
            print("Cannot create NDI sender.")
            return

        # Create a video frame
        video_frame = ndi.VideoFrameV2()

        while True:
            # Capture the screen
            screen = pyautogui.screenshot()
            screen = screen.resize((1920, 1080), Image.LANCZOS)
            video_frame.data = screen.tobytes()
            video_frame.line_stride_in_bytes = screen.width * 3
            video_frame.xres = screen.width
            video_frame.yres = screen.height

            # Send the frame
            ndi.send_send_video_v2(sender, video_frame)

    @staticmethod
    def get_local_ip():
        s = socket.socket(socket.AF_INET, socket.SOCK_DGRAM)
        try:
            s.connect(("10.254.254.254", 1))
            ip = s.getsockname()[0]
        except Exception:
            ip = "127.0.0.1"
        finally:
            s.close()
        return ip

@app.route('/get_ip', methods=['GET'])
def get_ip():
    ip = App.get_local_ip()
    return jsonify({'ip': ip})

@app.route('/next', methods=['GET'])
def next_slide():
    app_instance.next_slide()
    return jsonify({'status': 'ok'})

@app.route('/prev', methods=['GET'])
def prev_slide():
    app_instance.prev_slide()
    return jsonify({'status': 'ok'})

@app.route('/show_screen', methods=['GET'])
def show_screen():
    app_instance.show_screen()
    return jsonify({'status': 'ok'})

def run_flask_app():
    app.run(host='0.0.0.0', port=5000)

if __name__ == "__main__":
    # Start Flask server in a separate thread
    flask_thread = threading.Thread(target=run_flask_app)
    flask_thread.daemon = True
    flask_thread.start()

    root = tk.Tk()
    app_instance = App(root)
    root.mainloop()